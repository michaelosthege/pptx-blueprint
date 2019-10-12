import pathlib
import pptx
from typing import Union

__Pathlike = Union[str, pathlib.Path]


class TemplateOvewriteError(Exception):
    """An error that occurs when the user tries to override the template."""
    pass


class Template:
    """Helper class for modifying pptx templates.
    """

    def __init__(self, filename: __Pathlike):
        """Initializes a Template-Modifier.

        Args:
            filename (path-like): file name or path
        """
        self._template_path = filename
        self._presentation = pptx.Presentation(filename)
        pass

    def replace_text(self, label: str, text: str, *, scope=None):
        """Replaces text placeholders on one or many slides.

        Args:
            label (str): label of the placeholder (without curly braces)
            text (str): new content
            scope: None, slide number, Slide object or iterable of Slide objects
        """
        pass

    def replace_picture(self, label: str, filename: __Pathlike):
        """Replaces rectangle placeholders on one or many slides.

        Args:
            label (str): label of the placeholder (without curly braces)
            filename (path-like): path to an image file
        """
        pass

    def replace_table(self, label: str, data):
        """Replaces rectangle placeholders on one or many slides.

        Args:
            label (str): label of the placeholder (without curly braces)
            data (pandas.DataFrame): table to be inserted into the presentation
        """
        pass

    def save(self, filename: __Pathlike):
        """Saves the updated pptx to the specified filepath.

        Args:
            filename (path-like): file name or path

        Raises:
            TemplateOvewriteError
        """
        if pathlib.Path(filename).absolute() == pathlib.Path(self._template_path).absolute():
            raise TemplateOvewriteError(
                f'The specified save path ({filename}) is equal to the path of the template file.' \
                'The template should not be overwritten.'
            )
        self._presentation.save(filename)
